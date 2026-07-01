"""依「1142 GenAI 期末個人專題簡報模板」產生 AIR-Agent 專題簡報 PPTX。

作法：直接以模板為基底（沿用母片、版面配置、字型與配色），清空原示範投影片後，
依報告內容重建約 15 頁。字型採微軟正黑體、粗體；每頁保留頁碼欄位。

用法：uv run python scripts/build_slides.py
"""
import copy

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

TEMPLATE = "1142 GenAI 期末個人專題簡報模板.pptx"
OUT = "專題簡報_AIR-Agent.pptx"
FONT = "微軟正黑體"


def clear_slides(prs):
    """移除示範投影片，連同其 part 與關聯，避免序列化時 partname 衝突。"""
    part = prs.part
    sld_id_lst = prs.slides._sldIdLst
    for sid in list(sld_id_lst):
        rId = sid.get(qn("r:id"))
        part.drop_rel(rId)
        sld_id_lst.remove(sid)


def style_run(run, size=None, bold=True):
    run.font.name = FONT
    run.font.bold = bold
    if size is not None:
        run.font.size = Pt(size)
    # 套用東亞字型
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", FONT)


def set_title(slide, text, size=None):
    title = slide.shapes.title
    title.text = ""
    p = title.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    style_run(run, size=size, bold=True)


def get_body(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            return ph
    return None


def fill_bullets(slide, items):
    """items: list of (text, level) 或 純字串(level 0)。"""
    body = get_body(slide)
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for it in items:
        text, level = (it if isinstance(it, tuple) else (it, 0))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        run = p.add_run()
        run.text = text
        style_run(run, bold=True)


def add_content_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, title)
    fill_bullets(slide, items)
    add_slide_number_field(slide)
    return slide


def add_slide_number_field(slide):
    """python-pptx 預設不會把頁碼佔位複製進新投影片，這裡從版面配置複製過來，
    使其沿用母片的頁碼欄位（slidenum）與位置／格式。"""
    if any(ph.placeholder_format.idx == 4 for ph in slide.placeholders):
        return
    for ph in slide.slide_layout.placeholders:
        if ph.placeholder_format.idx == 4:
            slide.shapes._spTree.append(copy.deepcopy(ph._element))
            return


def add_table_slide(prs, title, header, rows, col_ratio=None):
    from pptx.util import Emu
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, title)
    body = get_body(slide)
    left, top, width, height = body.left, body.top, body.width, body.height
    # 移除內容佔位（用表格取代）
    body._element.getparent().remove(body._element)
    nrows, ncols = len(rows) + 1, len(header)
    gtable = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    table = gtable.table
    if col_ratio:
        total = sum(col_ratio)
        for i, r in enumerate(col_ratio):
            table.columns[i].width = Emu(int(width * r / total))
    for j, h in enumerate(header):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                style_run(r, size=24, bold=True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    style_run(r, size=18, bold=False)
    add_slide_number_field(slide)
    return slide


def add_title_slide(prs, title, subtitle_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(slide, title)
    # 副標題
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, line in enumerate(subtitle_lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = line
                style_run(run, bold=True)
    add_slide_number_field(slide)
    return slide


def build():
    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    # 1. 封面
    add_title_slide(
        prs,
        "AIR-Agent：具長期記憶與自主行為的\nAI 論文研究代理人系統設計",
        ["GenAI 期末個人專題", "黃士育", "2026 年 6 月"],
    )

    # 2. 大綱
    add_content_slide(prs, "大綱", [
        "專案背景與問題定義",
        "系統總覽與核心工作流",
        "功能展示：斜線指令",
        "架構設計：六大支柱",
        "工程取捨與應用情境",
        "限制、未來工作與結論",
    ])

    # 3. 背景與問題
    add_content_slide(prs, "專案背景與問題定義", [
        "痛點一｜資訊過載：arXiv 每日數百篇，人工追蹤成本高、英文閱讀門檻造成延遲",
        "痛點二｜知識零散：讀過的論文缺乏結構化沉澱，難以跨論文關聯查詢",
        "痛點三｜調研耗時：單一主題文獻回顧往往要花數小時",
        "問題定義：給定持續更新的論文來源與自然語言意圖，系統應自主決定",
        ("抓什麼、如何濃縮、如何沉澱為可檢索知識，並基於可信來源生成有引用依據的回答", 1),
        "本質：RAG ＋ 排程式自主行為的代理人問題",
    ])

    # 4. 系統總覽
    add_content_slide(prs, "系統總覽：分層架構", [
        "介面層：Discord 斜線指令 /daily /ask /report /set_push_time /help",
        "控制層：指令路由、tasks.loop 排程、狀態持久化（schedule.json）",
        "插件層：arXiv 爬蟲、GroqClient（LLM）、FAISS 向量庫＋本地嵌入器",
        "資源層：arXiv API、Groq API、本地 sentence-transformers 模型",
        "核心模型：Groq 託管 llama-3.3-70b-versatile（OpenAI 相容、低延遲 LPU）",
    ])

    # 5. 三條工作流（表格）
    add_table_slide(prs, "三條核心工作流",
                    ["工作流", "觸發", "行為"],
                    [["每日情報推送", "排程（每日）", "抓最新論文→逐篇摘要→推送 Embed→寫入向量庫"],
                     ["RAG 問答", "/ask <問題>", "向量檢索 Top-k→帶來源生成答案→附參考論文"],
                     ["主題研究報告", "/report <主題>", "自主檢索→多篇證據彙整→結構化報告→沉澱知識"]],
                    col_ratio=[2, 2, 6])

    # 6. 功能展示：斜線指令
    add_table_slide(prs, "功能展示：斜線指令一覽",
                    ["指令", "功能"],
                    [["/daily", "立即抓取並推送今日 AI 論文"],
                     ["/ask <問題>", "依向量庫論文做 RAG 問答，附參考論文"],
                     ["/report <主題>", "自主檢索相關論文並產出結構化研究報告"],
                     ["/set_push_time <時> <分>", "設定每日自動推送時間（即時生效並持久化）"],
                     ["/help", "顯示指令說明與目前推送時間"]],
                    col_ratio=[4, 6])

    # 7. System Prompt
    add_content_slide(prs, "架構支柱 1：System Prompt 設計", [
        "核心理念：分任務專屬 prompt，而非單一萬用 prompt → 品質與成本雙控",
        "摘要：明確角色／語言、長度約束（2-3 句）、聚焦問題與方法、去客套",
        ("→ 解決跨語言問題，並直接壓低輸出 token", 1),
        "問答：強制「僅根據來源」grounding、誠實退場、列出參考論文",
        ("→ 抑制幻覺、可溯源，符合學術剛性需求", 1),
        "報告：prompt 鎖定固定 Markdown 結構 → 下游顯示／轉檔可預測可解析",
    ])

    # 7. LLM 選型
    add_content_slide(prs, "架構支柱 2：LLM 模型選擇", [
        "能力：摘要與多篇彙整屬中等難度，70B 開源模型已足夠",
        "延遲：Discord 互動有 3 秒回應限制 → Groq LPU 低延遲為決定性因素",
        "成本：開源模型定價低，適合高頻呼叫",
        "生態：Groq 提供 OpenAI 相容端點，沿用既有 SDK、零學習成本",
        "為何不用旗艦模型？邊際品質提升無法正當化數十倍單價",
        "LLM 封裝於 GroqClient，模型由設定檔注入 → 可替換不影響工作流",
    ])

    # 8. Knowledge
    add_content_slide(prs, "架構支柱 3：Knowledge 知識庫", [
        "佈局：來源 → 濃縮 → 向量化 → 持久化",
        "Chunk 策略：一篇論文＝一個 chunk（標題＋摘要），不再切段以保語意完整",
        "Metadata 與向量分離：FAISS 存向量、metadata.json 存原始欄位，索引一一對應",
        "相似度：L2 正規化後內積＝餘弦相似度 → IndexFlatIP 精確檢索、零召回損失",
        "去重：以 arXiv 短 ID 為主鍵，加入前先去重，避免知識庫膨脹",
    ])

    # 9. Plugin
    add_content_slide(prs, "架構支柱 4：Plugin 插件設計", [
        "插件＝代理人可調用的能力單元（tools），各有清晰輸入／輸出契約",
        "ArxivCrawler：fetch_latest_papers / search_topic",
        "GroqClient：summarize / answer / research_report",
        "VectorStore：add 收錄、search 檢索；Embedder：encode 向量化",
        "設計原則：單一職責＋介面穩定＋可替換（來源／LLM／向量庫皆可換）",
    ])

    # 10. Memory
    add_content_slide(prs, "架構支柱 5：Memory 記憶機制", [
        "三層記憶設計：",
        ("長期語意記憶｜FAISS＋metadata.json｜永久｜沉澱論文支撐 RAG", 1),
        ("狀態記憶｜schedule.json＋去重集合｜永久｜記住推送時間與已收錄", 1),
        ("工作記憶｜context window｜單次｜把 Top-k 證據塞入 prompt", 1),
        "關鍵：以本地向量庫實現長期記憶，跨重啟存活、可累積、可去重",
        "RAG 本質＝生成時從長期記憶動態載入最相關片段到工作記憶",
    ])

    # 11. Workflow
    add_content_slide(prs, "架構支柱 6：Workflow 工作流控制", [
        "非阻塞執行：爬蟲與 LLM 呼叫以 asyncio.to_thread 丟背景，不阻塞事件迴圈",
        "逾時防護：互動指令先 defer()，將 3 秒硬限制延展為 15 分鐘 followup",
        "容錯不中斷：找不到頻道／LLM 失敗皆有 try/except 與降級回覆",
        "動態重排程：/set_push_time 透過 change_interval() 即時生效、無需重啟",
    ])

    # 12. 工程取捨與容錯
    add_content_slide(prs, "工程取捨與容錯設計", [
        "幻覺抑制：grounding ＋ 強制引用，回答可溯源、符合學術需求",
        "精確檢索：資料規模小，選 IndexFlatIP 精確檢索、零召回損失",
        "本地嵌入：免費、離線、保護隱私，不依賴外部嵌入服務",
        "金鑰管理：Groq／Discord token 以 .env 注入、不進版控",
        "容錯不中斷：找不到頻道／LLM 失敗皆降級回覆，單點失敗不致服務崩潰",
    ])

    # 13. 應用情境與價值
    add_content_slide(prs, "應用情境與價值", [
        "個人研究者：每日自動掌握 arXiv 最新進展，免去人工追蹤",
        "研究社群／實驗室：共享 Discord 頻道，集體沉澱知識庫",
        "文獻回顧加速：/report 主題一鍵產出結構化研究報告與引用",
        "跨語言友善：中文提問、英文文獻檢索，降低閱讀門檻",
        "越用越聰明：每次互動都把論文沉澱進長期記憶，知識持續累積",
    ])

    # 15. 限制與未來
    add_content_slide(prs, "限制與未來工作", [
        "限制：來源單一（僅 arXiv）；純向量檢索對專名匹配較弱；受 70B 模型上限約束",
        "未來：多來源整合（HuggingFace、會議論文、GitHub Trending）",
        "混合檢索 Hybrid Search（dense＋BM25）提升專名召回",
        "Re-ranking：cross-encoder 重排提升 Top-k 精準度",
        "多語言嵌入升級、Prompt／語意快取、主題演化趨勢分析",
    ])

    # 16. 結論
    add_content_slide(prs, "結論", [
        "完整閉環：自動情報蒐集 → 知識沉澱 → RAG 問答 → 主題報告",
        "架構面：六大支柱各司其職、鬆耦合、可替換",
        "功能面：斜線指令直覺易用，涵蓋推送、問答、主題報告與排程設定",
        "技術面：跨語言檢索、Agentic 工作流、長期記憶沉澱、容錯不中斷",
        "結語：一個具備長期記憶與自主行為的可運作研究代理人",
    ])

    prs.save(OUT)
    print(f"已輸出：{OUT}（{len(prs.slides._sldIdLst)} 頁）")


if __name__ == "__main__":
    build()
